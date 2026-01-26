"""
Gmail Watcher V2 - Enhanced Edition
- Continuously monitors Gmail for new emails
- Extracts text from 20+ file types
- Improved CSV/Excel/Image handling
- Real-time indexing with Pathway
"""

import os
import base64
import pickle
import time
from pathlib import Path
from datetime import datetime
from dateutil import parser
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import logging

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Gmail API scopes
SCOPES = ['https://www.googleapis.com/auth/gmail.readonly']

# Configuration
OUTPUT_DIR = Path("data/emails")
ATTACHMENTS_DIR = Path("data/attachments")
PROCESSED_DIR = Path("data/emails_processed")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
ATTACHMENTS_DIR.mkdir(parents=True, exist_ok=True)
PROCESSED_DIR.mkdir(parents=True, exist_ok=True)

# Track which emails we've already downloaded
TRACKING_FILE = Path("gmail_sync_state.txt")

# ===== ENHANCED ATTACHMENT CONFIGURATION =====
ATTACHMENT_CONFIG = {
    'enabled': True,
    'max_file_size_mb': 25,  # Increased limit
    'allowed_extensions': [
        # Documents
        '.pdf', '.docx', '.doc', '.txt', '.rtf', '.odt',
        # Spreadsheets
        '.csv', '.xlsx', '.xls', '.ods',
        # Presentations
        '.pptx', '.ppt', '.odp',
        # Images
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',
        # Web/Markup
        '.html', '.htm', '.xml', '.json', '.md',
        # Code/Scripts
        '.py', '.js', '.java', '.cpp', '.c', '.sh', '.sql',
        # Archives (metadata only)
        '.zip', '.rar', '.7z', '.tar', '.gz',
        # Other
        '.log', '.yaml', '.yml', '.ini', '.conf'
    ],
    'pdf_max_pages': 50,  # Increased limit
    'save_metadata_only_if_too_large': True,
    'extract_text_immediately': True,
}


# ===== ENHANCED TEXT EXTRACTION FUNCTIONS =====

def extract_pdf_text(file_path: Path) -> str:
    """Extract text from PDF with better error handling"""
    try:
        import pypdf
        text_parts = []
        with open(file_path, 'rb') as f:
            pdf_reader = pypdf.PdfReader(f)
            total_pages = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"--- Page {page_num}/{total_pages} ---\n{page_text.strip()}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num}: {e}")
                    continue
        
        if text_parts:
            return "\n\n".join(text_parts)
        return "PDF contains no extractable text (may be scanned images)"
    
    except ImportError:
        logger.error("PyPDF2 not installed. Install: pip install PyPDF2")
        return None
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return None


def extract_docx_text(file_path: Path) -> str:
    """Extract text from DOCX with tables"""
    try:
        import docx
        doc = docx.Document(file_path)
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables, 1):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                text_parts.append(f"\n--- Table {table_idx} ---")
                if len(table_data) > 0:
                    headers = table_data[0]
                    text_parts.append("| " + " | ".join(headers) + " |")
                    text_parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    for row in table_data[1:]:
                        while len(row) < len(headers):
                            row.append("")
                        text_parts.append("| " + " | ".join(row[:len(headers)]) + " |")
        
        return "\n".join(text_parts) if text_parts else "Document is empty"
    
    except ImportError:
        logger.error("python-docx not installed. Install: pip install python-docx")
        return None
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return None


def extract_doc_text(file_path: Path) -> str:
    """Extract text from legacy DOC files"""
    try:
        import textract
        text = textract.process(str(file_path)).decode('utf-8')
        return f"Legacy Word Document (.doc)\n\n{text}" if text.strip() else "Document is empty"
    except ImportError:
        logger.warning("textract not installed. Install: pip install textract")
        return "Legacy .doc file - install 'textract' to extract content"
    except Exception as e:
        logger.error(f"Error extracting DOC: {e}")
        return None


def extract_rtf_text(file_path: Path) -> str:
    """Extract text from RTF files"""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        return f"RTF Document\n\n{text}" if text.strip() else "Document is empty"
    except ImportError:
        logger.warning("striprtf not installed. Install: pip install striprtf")
        return None
    except Exception as e:
        logger.error(f"Error extracting RTF: {e}")
        return None


def extract_odt_text(file_path: Path) -> str:
    """Extract text from OpenDocument Text files"""
    try:
        from odf import text, teletype
        from odf.opendocument import load
        
        textdoc = load(str(file_path))
        allparas = textdoc.getElementsByType(text.P)
        text_parts = [teletype.extractText(para) for para in allparas]
        
        return f"OpenDocument Text\n\n" + "\n".join(text_parts) if text_parts else "Document is empty"
    except ImportError:
        logger.warning("odfpy not installed. Install: pip install odfpy")
        return None
    except Exception as e:
        logger.error(f"Error extracting ODT: {e}")
        return None


def _manual_markdown_table(df) -> str:
    """Manual markdown table creation"""
    try:
        lines = []
        headers = [str(col)[:50] for col in df.columns]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        
        for _, row in df.iterrows():
            row_values = [str(val)[:100] if val != "" else " " for val in row.values]
            lines.append("| " + " | ".join(row_values) + " |")
        
        return "\n".join(lines)
    except Exception as e:
        return str(df)


def extract_csv_text(file_path: Path) -> str:
    """Extract text from CSV with robust parsing"""
    try:
        import pandas as pd
        
        df = None
        encoding_used = None
        sep_used = ','
        
        # Try multiple encodings with comma separator
        encodings = ['utf-8', 'latin1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                df = pd.read_csv(
                    file_path, 
                    encoding=encoding,
                    on_bad_lines='skip',
                    engine='python'
                )
                encoding_used = encoding
                break
            except:
                continue
        
        # Try different separators if comma failed
        if df is None or df.empty:
            separators = [';', '\t', '|']
            for sep in separators:
                for encoding in encodings:
                    try:
                        df = pd.read_csv(
                            file_path,
                            sep=sep,
                            encoding=encoding,
                            on_bad_lines='skip',
                            engine='python'
                        )
                        if not df.empty:
                            encoding_used = encoding
                            sep_used = sep
                            break
                    except:
                        continue
                if df is not None and not df.empty:
                    break
        
        if df is None or df.empty:
            return f"Could not parse CSV: {file_path.name}"
        
        # Clean data
        df = df.fillna("")
        df = df.loc[:, df.any()]
        df = df[df.any(axis=1)]
        
        # Strip whitespace
        df.columns = df.columns.astype(str).str.strip()
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].astype(str).str.strip()
        
        text_parts = [
            f"CSV File: {file_path.name}",
            f"Encoding: {encoding_used}, Separator: '{sep_used}'",
            f"Rows: {len(df)}, Columns: {len(df.columns)}",
            f"Columns: {', '.join(df.columns.astype(str))}",
            "\n--- Data (First 500 rows) ---\n"
        ]
        
        preview_df = df.head(500)
        
        try:
            table_md = preview_df.to_markdown(index=False, tablefmt="github")
        except:
            table_md = _manual_markdown_table(preview_df)
        
        text_parts.append(table_md)
        
        if len(df) > 500:
            text_parts.append(f"\n(Showing first 500 of {len(df)} rows)")
        
        # Summary statistics
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            text_parts.append("\n--- Numeric Statistics ---")
            for col in numeric_cols[:10]:
                try:
                    stats = df[col].describe()
                    text_parts.append(
                        f"{col}: Min={stats['min']:.2f}, Max={stats['max']:.2f}, Mean={stats['mean']:.2f}"
                    )
                except:
                    pass
        
        return "\n".join(text_parts)
    
    except ImportError:
        logger.error("pandas not installed. Install: pip install pandas")
        return None
    except Exception as e:
        logger.error(f"Error extracting CSV: {e}")
        return f"Error parsing CSV: {e}"


def extract_excel_text(file_path: Path) -> str:
    """Extract text from Excel with SMART header detection"""
    try:
        import pandas as pd
        import openpyxl
        
        text_parts = [f"Excel File: {file_path.name}\n"]
        
        # Load workbook with openpyxl to handle merged cells
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
            ws = wb[sheet_name]
            
            # Convert to list of lists
            data = []
            for row in ws.iter_rows(values_only=True):
                # Skip completely empty rows
                if any(cell for cell in row if cell):
                    data.append(list(row))
            
            if not data:
                text_parts.append(f"\n{'='*60}\nSheet {sheet_idx}: {sheet_name}\n(Empty)\n")
                continue
            
            # Smart header detection
            header_row_idx = None
            
            # Look for row with day names (common in menus)
            day_keywords = ['monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday']
            for idx, row in enumerate(data[:10]):  # Check first 10 rows
                row_str = ' '.join([str(cell).lower() for cell in row if cell])
                if sum(day in row_str for day in day_keywords) >= 3:  # At least 3 days found
                    header_row_idx = idx
                    break
            
            # If no day names, look for row with multiple non-empty strings
            if header_row_idx is None:
                for idx, row in enumerate(data[:5]):
                    non_empty = [cell for cell in row if cell and isinstance(cell, str)]
                    if len(non_empty) >= 3:
                        header_row_idx = idx
                        break
            
            # Build structured text
            text_parts.append(f"\n{'='*60}")
            text_parts.append(f"Sheet {sheet_idx}: {sheet_name}")
            
            # Add title rows (before header)
            if header_row_idx and header_row_idx > 0:
                text_parts.append("\n--- Title/Description ---")
                for row in data[:header_row_idx]:
                    row_text = ' | '.join([str(cell) for cell in row if cell])
                    if row_text:
                        text_parts.append(row_text)
            
            # Create DataFrame with proper headers
            if header_row_idx is not None:
                headers = data[header_row_idx]
                content_data = data[header_row_idx + 1:]
                
                # Clean headers
                clean_headers = []
                for i, h in enumerate(headers):
                    if h and str(h).strip():
                        clean_headers.append(str(h).strip())
                    else:
                        clean_headers.append(f"Column_{i+1}")
                
                # Create DataFrame
                df = pd.DataFrame(content_data, columns=clean_headers)
                
                # Remove empty columns
                df = df.loc[:, df.any()]
                
                # Remove completely empty rows
                df = df[df.any(axis=1)]
                
                # Replace NaN with empty string
                df = df.fillna("")
                
                # Clean up asterisks and empty values
                for col in df.columns:
                    if df[col].dtype == 'object':
                        df[col] = df[col].apply(lambda x: "" if str(x).strip() in ['***', '****', 'nan', 'None'] else str(x).strip())
                
                text_parts.append(f"\nRows: {len(df)}, Columns: {len(df.columns)}")
                text_parts.append(f"Headers: {', '.join(clean_headers)}")
                text_parts.append("\n--- Structured Data ---\n")
                
                # IMPROVED TABLE FORMATTING
                # Group by meal type (detect sections)
                current_section = None
                section_data = []
                
                for idx, row in df.iterrows():
                    # Detect section headers (BREAKFAST, LUNCH, etc.)
                    first_col = str(row.iloc[0]).upper()
                    
                    if any(keyword in first_col for keyword in ['BREAK FAST', 'BREAKFAST', 'LUNCH', 'DINNER', 'SNACK', 'EVENING']):
                        # Save previous section
                        if current_section and section_data:
                            text_parts.append(f"\n### {current_section}")
                            text_parts.append(format_menu_section(section_data, clean_headers))
                        
                        current_section = first_col
                        section_data = []
                    else:
                        section_data.append(row)
                
                # Add last section
                if current_section and section_data:
                    text_parts.append(f"\n### {current_section}")
                    text_parts.append(format_menu_section(section_data, clean_headers))
                
            else:
                # Fallback: no header detected
                text_parts.append("\n--- Raw Data ---")
                for row in data[:100]:  # Limit to 100 rows
                    row_text = ' | '.join([str(cell) for cell in row if cell])
                    if row_text:
                        text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    except ImportError:
        logger.error("pandas/openpyxl not installed. Install: pip install pandas openpyxl")
        return None
    except Exception as e:
        logger.error(f"Error extracting Excel: {e}")
        return f"Error reading Excel: {e}"


def format_menu_section(section_data, headers):
    """Format menu section in a readable way"""
    lines = []
    
    # Transpose the data for better readability (days as sections)
    for col_idx, header in enumerate(headers):
        if header.startswith('Column_'):
            continue
            
        lines.append(f"\n**{header}:**")
        
        for row in section_data:
            item = str(row.iloc[col_idx]) if col_idx < len(row) else ""
            if item and item.strip():
                lines.append(f"  • {item}")
    
    return '\n'.join(lines)


def extract_ods_text(file_path: Path) -> str:
    """Extract text from OpenDocument Spreadsheet"""
    try:
        import pandas as pd
        df = pd.read_excel(file_path, engine='odf')
        
        text_parts = [
            f"OpenDocument Spreadsheet: {file_path.name}",
            f"Rows: {len(df)}, Columns: {len(df.columns)}",
            "\n--- Data ---\n"
        ]
        
        try:
            table_md = df.head(200).to_markdown(index=False, tablefmt="github")
        except:
            table_md = _manual_markdown_table(df.head(200))
        
        text_parts.append(table_md)
        return "\n".join(text_parts)
    
    except ImportError:
        logger.warning("odfpy not installed. Install: pip install odfpy")
        return None
    except Exception as e:
        logger.error(f"Error extracting ODS: {e}")
        return None


def extract_pptx_text(file_path: Path) -> str:
    """Extract text from PowerPoint presentations"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = [f"PowerPoint Presentation: {file_path.name}\n"]
        text_parts.append(f"Total Slides: {len(prs.slides)}\n")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_parts.append(f"\n--- Slide {slide_idx} ---")
                text_parts.append("\n".join(slide_text))
        
        return "\n".join(text_parts) if len(text_parts) > 2 else "Presentation is empty"
    
    except ImportError:
        logger.warning("python-pptx not installed. Install: pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return None


def extract_image_text(file_path: Path) -> str:
    """Enhanced OCR extraction with preprocessing"""
    try:
        from PIL import Image, ImageEnhance, ImageFilter
        import pytesseract
        
        image = Image.open(file_path)
        
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Preprocessing
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        image = image.filter(ImageFilter.SHARPEN)
        image = image.convert('L')
        
        # Try multiple OCR configs
        configs = ['--psm 6', '--psm 3', '--psm 11']
        best_text = ""
        
        for config in configs:
            try:
                text = pytesseract.image_to_string(image, config=config)
                if len(text.strip()) > len(best_text.strip()):
                    best_text = text
            except:
                continue
        
        if best_text.strip():
            return f"Image OCR: {file_path.name}\n\n--- Extracted Text ---\n{best_text.strip()}"
        else:
            return f"Image '{file_path.name}' contains no readable text"
    
    except ImportError:
        logger.error("PIL/pytesseract not installed. Install: pip install pillow pytesseract")
        return None
    except Exception as e:
        logger.error(f"Error extracting image text: {e}")
        return None


def extract_html_text(file_path: Path) -> str:
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style tags
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text(separator='\n', strip=True)
        return f"HTML Document: {file_path.name}\n\n{text}" if text else "HTML is empty"
    
    except ImportError:
        logger.warning("beautifulsoup4 not installed. Install: pip install beautifulsoup4")
        return None
    except Exception as e:
        logger.error(f"Error extracting HTML: {e}")
        return None


def extract_xml_text(file_path: Path) -> str:
    """Extract text from XML files"""
    try:
        import xml.etree.ElementTree as ET
        
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def extract_text_recursive(element, level=0):
            lines = []
            indent = "  " * level
            
            if element.text and element.text.strip():
                lines.append(f"{indent}{element.tag}: {element.text.strip()}")
            else:
                lines.append(f"{indent}{element.tag}")
            
            for child in element:
                lines.extend(extract_text_recursive(child, level + 1))
            
            return lines
        
        text_lines = extract_text_recursive(root)
        return f"XML Document: {file_path.name}\n\n" + "\n".join(text_lines)
    
    except Exception as e:
        logger.error(f"Error extracting XML: {e}")
        return None


def extract_json_text(file_path: Path) -> str:
    """Extract and format JSON files"""
    try:
        import json
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        return f"JSON Document: {file_path.name}\n\n{formatted}"
    
    except Exception as e:
        logger.error(f"Error extracting JSON: {e}")
        return None


def extract_markdown_text(file_path: Path) -> str:
    """Extract text from Markdown files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return f"Markdown Document: {file_path.name}\n\n{content}"
    except Exception as e:
        logger.error(f"Error reading Markdown: {e}")
        return None


def extract_code_text(file_path: Path) -> str:
    """Extract text from code files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        ext = file_path.suffix
        lang_map = {
            '.py': 'Python', '.js': 'JavaScript', '.java': 'Java',
            '.cpp': 'C++', '.c': 'C', '.sh': 'Shell', '.sql': 'SQL',
            '.yaml': 'YAML', '.yml': 'YAML', '.ini': 'INI', '.conf': 'Config'
        }
        
        lang = lang_map.get(ext, 'Code')
        return f"{lang} File: {file_path.name}\n\n```{ext[1:]}\n{content}\n```"
    
    except Exception as e:
        logger.error(f"Error reading code file: {e}")
        return None


def extract_text_file(file_path: Path) -> str:
    """Extract text from plain text files"""
    try:
        encodings = ['utf-8', 'latin1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                    content = f.read()
                    if content.strip():
                        return f"Text File: {file_path.name}\n\n{content}"
            except:
                continue
        
        return f"Could not read: {file_path.name}"
    
    except Exception as e:
        logger.error(f"Error reading text file: {e}")
        return None


def extract_archive_metadata(file_path: Path) -> str:
    """Extract metadata from archive files"""
    try:
        import zipfile
        import rarfile
        import tarfile
        
        ext = file_path.suffix.lower()
        files_list = []
        
        if ext == '.zip':
            with zipfile.ZipFile(file_path, 'r') as zf:
                files_list = [f.filename for f in zf.filelist]
        elif ext == '.rar':
            with rarfile.RarFile(file_path, 'r') as rf:
                files_list = rf.namelist()
        elif ext in ['.tar', '.gz', '.tgz']:
            with tarfile.open(file_path, 'r:*') as tf:
                files_list = tf.getnames()
        
        if files_list:
            return f"Archive: {file_path.name}\n\nContents ({len(files_list)} files):\n" + "\n".join(files_list[:100])
        
        return f"Archive: {file_path.name} (could not list contents)"
    
    except Exception as e:
        logger.error(f"Error reading archive: {e}")
        return f"Archive: {file_path.name} (metadata extraction failed)"


def extract_text_from_file(file_path: Path) -> str:
    """Master function to extract text from any supported file type"""
    ext = file_path.suffix.lower()
    
    extractors = {
        '.pdf': extract_pdf_text,
        '.docx': extract_docx_text,
        '.doc': extract_doc_text,
        '.rtf': extract_rtf_text,
        '.odt': extract_odt_text,
        '.csv': extract_csv_text,
        '.xlsx': extract_excel_text,
        '.xls': extract_excel_text,
        '.ods': extract_ods_text,
        '.pptx': extract_pptx_text,
        '.ppt': extract_pptx_text,
        '.jpg': extract_image_text,
        '.jpeg': extract_image_text,
        '.png': extract_image_text,
        '.gif': extract_image_text,
        '.bmp': extract_image_text,
        '.tiff': extract_image_text,
        '.webp': extract_image_text,
        '.html': extract_html_text,
        '.htm': extract_html_text,
        '.xml': extract_xml_text,
        '.json': extract_json_text,
        '.md': extract_markdown_text,
        '.py': extract_code_text,
        '.js': extract_code_text,
        '.java': extract_code_text,
        '.cpp': extract_code_text,
        '.c': extract_code_text,
        '.sh': extract_code_text,
        '.sql': extract_code_text,
        '.yaml': extract_code_text,
        '.yml': extract_code_text,
        '.ini': extract_code_text,
        '.conf': extract_code_text,
        '.log': extract_text_file,
        '.txt': extract_text_file,
        '.zip': extract_archive_metadata,
        '.rar': extract_archive_metadata,
        '.7z': extract_archive_metadata,
        '.tar': extract_archive_metadata,
        '.gz': extract_archive_metadata,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    
    return None


# ===== GMAIL FUNCTIONS =====

def authenticate_gmail():
    """Authenticate with Gmail API"""
    creds = None
    
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not os.path.exists('credentials.json'):
                print("\n❌ credentials.json not found!")
                print("Run: python fetch_gmail.py first to set up authentication")
                return None
            
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)
    
    return creds


def load_synced_ids():
    """Load set of already-synced email IDs"""
    if TRACKING_FILE.exists():
        with open(TRACKING_FILE, 'r') as f:
            return set(line.strip() for line in f if line.strip())
    return set()


def save_synced_id(email_id):
    """Save a synced email ID"""
    with open(TRACKING_FILE, 'a') as f:
        f.write(f"{email_id}\n")


def get_email_body(payload):
    """Extract email body from message payload - supports both plain text and HTML"""
    plain_body = ""
    html_body = ""
    
    def extract_body_recursive(parts):
        nonlocal plain_body, html_body
        
        for part in parts:
            mime_type = part.get('mimeType', '')
            
            # Check for nested parts
            if 'parts' in part:
                extract_body_recursive(part['parts'])
            
            # Extract plain text
            if mime_type == 'text/plain' and 'data' in part.get('body', {}):
                plain_body = base64.urlsafe_b64decode(part['body']['data']).decode('utf-8', errors='ignore')
            
            # Extract HTML
            elif mime_type == 'text/html' and 'data' in part.get('body', {}):
                html_body = base64.urlsafe_b64decode(part['body']['data']).decode('utf-8', errors='ignore')
    
    # Handle multipart emails
    if 'parts' in payload:
        extract_body_recursive(payload['parts'])
    # Handle simple emails
    elif 'body' in payload and 'data' in payload['body']:
        mime_type = payload.get('mimeType', '')
        body_data = base64.urlsafe_b64decode(payload['body']['data']).decode('utf-8', errors='ignore')
        
        if mime_type == 'text/html':
            html_body = body_data
        else:
            plain_body = body_data
    
    # Prefer plain text, but use HTML if that's all we have
    if plain_body:
        return plain_body
    elif html_body:
        # Convert HTML to plain text
        return html_to_text(html_body)
    
    return ""


def html_to_text(html_content):
    """Convert HTML email to plain text"""
    try:
        from bs4 import BeautifulSoup
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style", "head", "meta", "link"]):
            script.decompose()
        
        # Get text
        text = soup.get_text(separator='\n', strip=True)
        
        # Clean up multiple newlines
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        
        return '\n'.join(lines)
    
    except ImportError:
        logger.warning("beautifulsoup4 not installed - falling back to raw HTML")
        # Basic fallback: strip HTML tags
        import re
        text = re.sub('<[^<]+?>', '', html_content)
        return text.strip()
    
    except Exception as e:
        logger.error(f"Error converting HTML to text: {e}")
        return html_content


def get_pdf_page_count(file_path):
    """Get number of pages in a PDF"""
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            return len(pdf_reader.pages)
    except:
        return None


def save_attachment_metadata(attachment_info, email_id):
    """Save metadata file for skipped attachments"""
    metadata_filename = f"{attachment_info['filename']}.meta"
    metadata_path = ATTACHMENTS_DIR / f"email_{email_id}" / metadata_filename
    metadata_path.parent.mkdir(parents=True, exist_ok=True)
    
    metadata_content = f"""Attachment Metadata (File Too Large)
Email ID: {email_id}
Filename: {attachment_info['filename']}
Size: {attachment_info['size_mb']:.2f} MB
MIME Type: {attachment_info['mime_type']}
Reason: {attachment_info['skip_reason']}
"""
    
    with open(metadata_path, 'w', encoding='utf-8') as f:
        f.write(metadata_content)


def process_attachments(service, message_id, payload, email_subject):
    """Process and download attachments - ENHANCED VERSION"""
    
    if not ATTACHMENT_CONFIG['enabled']:
        return []
    
    attachments_info = []
    
    def extract_attachments(parts, message_id):
        for part in parts:
            if 'parts' in part:
                extract_attachments(part['parts'], message_id)
            
            filename = part.get('filename', '')
            if filename:
                file_ext = Path(filename).suffix.lower()
                
                if file_ext not in ATTACHMENT_CONFIG['allowed_extensions']:
                    logger.debug(f"Skipped (type not allowed): {filename}")
                    continue
                
                attachment_id = part['body'].get('attachmentId')
                size_bytes = part['body'].get('size', 0)
                size_mb = size_bytes / (1024 * 1024)
                
                if size_mb > ATTACHMENT_CONFIG['max_file_size_mb']:
                    logger.info(f"Skipped (too large: {size_mb:.2f}MB): {filename}")
                    
                    if ATTACHMENT_CONFIG['save_metadata_only_if_too_large']:
                        attachment_info = {
                            'filename': filename,
                            'size_mb': size_mb,
                            'mime_type': part.get('mimeType', 'unknown'),
                            'attachment_id': attachment_id,
                            'skip_reason': f'Exceeds size limit ({ATTACHMENT_CONFIG["max_file_size_mb"]}MB)'
                        }
                        save_attachment_metadata(attachment_info, message_id)
                    continue
                
                if attachment_id:
                    try:
                        attachment = service.users().messages().attachments().get(
                            userId='me',
                            messageId=message_id,
                            id=attachment_id
                        ).execute()
                        
                        file_data = base64.urlsafe_b64decode(attachment['data'])
                        
                        email_attachment_dir = ATTACHMENTS_DIR / f"email_{message_id}"
                        email_attachment_dir.mkdir(parents=True, exist_ok=True)
                        
                        safe_filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-')).strip()
                        attachment_path = email_attachment_dir / safe_filename
                        
                        with open(attachment_path, 'wb') as f:
                            f.write(file_data)
                        
                        # Check PDF page count
                        skip_pdf = False
                        if file_ext == '.pdf' and ATTACHMENT_CONFIG['pdf_max_pages']:
                            page_count = get_pdf_page_count(attachment_path)
                            if page_count and page_count > ATTACHMENT_CONFIG['pdf_max_pages']:
                                logger.info(f"PDF too long ({page_count} pages): {safe_filename}")
                                
                                if ATTACHMENT_CONFIG['save_metadata_only_if_too_large']:
                                    save_attachment_metadata({
                                        'filename': safe_filename,
                                        'size_mb': size_mb,
                                        'mime_type': part.get('mimeType', 'unknown'),
                                        'attachment_id': attachment_id,
                                        'skip_reason': f'PDF has {page_count} pages (limit: {ATTACHMENT_CONFIG["pdf_max_pages"]})'
                                    }, message_id)
                                
                                attachment_path.unlink()
                                skip_pdf = True
                        
                        if not skip_pdf:
                            # EXTRACT TEXT IMMEDIATELY
                            extracted_text = None
                            if ATTACHMENT_CONFIG['extract_text_immediately']:
                                extracted_text = extract_text_from_file(attachment_path)
                                
                                if extracted_text:
                                    output_filename = f"attachment_{message_id}_{Path(safe_filename).stem}.txt"
                                    output_path = PROCESSED_DIR / output_filename
                                    
                                    content = f"""Attachment from Email: {message_id}
Original Filename: {safe_filename}
File Type: {file_ext}
File Path: {attachment_path}

--- Extracted Content ---

{extracted_text}
"""
                                    
                                    with open(output_path, 'w', encoding='utf-8') as f:
                                        f.write(content)
                                    
                                    logger.info(f"✅ Extracted: {output_filename}")
                            
                            attachments_info.append({
                                'filename': safe_filename,
                                'path': str(attachment_path),
                                'size_mb': size_mb,
                                'type': file_ext,
                                'text_extracted': extracted_text is not None
                            })
                            logger.info(f"📎 Downloaded ({size_mb:.2f}MB): {safe_filename}")
                    
                    except Exception as e:
                        logger.error(f"Error downloading {filename}: {e}")
    
    if 'parts' in payload:
        extract_attachments(payload['parts'], message_id)
    
    return attachments_info


def fetch_new_emails(service, synced_ids, search_query="", max_check=200):
    """Fetch only new emails that haven't been synced yet"""
    
    try:
        results = service.users().messages().list(
            userId='me',
            maxResults=max_check,
            q=search_query
        ).execute()
        
        messages = results.get('messages', [])
        
        if not messages:
            return 0
        
        new_messages = [msg for msg in messages if msg['id'] not in synced_ids]
        
        if not new_messages:
            return 0
        
        logger.info(f"📥 Found {len(new_messages)} new emails")
        
        saved_count = 0
        
        for message in new_messages:
            try:
                filename = f"email_{message['id']}.txt"
                filepath = OUTPUT_DIR / filename
                
                if filepath.exists():
                    logger.debug(f"Skipped: {filename} (already exists)")
                    save_synced_id(message['id'])
                    synced_ids.add(message['id'])
                    continue
                
                msg = service.users().messages().get(
                    userId='me',
                    id=message['id'],
                    format='full'
                ).execute()
                
                headers = msg['payload']['headers']
                subject = next((h['value'] for h in headers if h['name'] == 'Subject'), 'No Subject')
                sender = next((h['value'] for h in headers if h['name'] == 'From'), 'Unknown')
                date = next((h['value'] for h in headers if h['name'] == 'Date'), '')
                try:
                    # Convert to ISO 8601 so Pathway/LLM can sort/filter easily
                    dt = parser.parse(date)
                    iso_date = dt.strftime("%Y-%m-%d %H:%M:%S")
                except:
                    iso_date = date
                
                to = next((h['value'] for h in headers if h['name'] == 'To'), '')
                
                body = get_email_body(msg['payload'])
                
                if not body:
                    save_synced_id(message['id'])
                    continue
                
                logger.info(f"📧 Processing: {subject[:60]}...")
                attachments = process_attachments(service, message['id'], msg['payload'], subject)
                
                # Extract additional metadata
                cc = next((h['value'] for h in headers if h['name'] == 'Cc'), '')
                reply_to = next((h['value'] for h in headers if h['name'] == 'Reply-To'), '')
                
                # Get labels (categories)
                labels = msg.get('labelIds', [])
                label_str = ', '.join(labels) if labels else 'None'
                
                # Create enhanced email content with metadata
                email_content = f"""From: {sender}
To: {to}
Cc: {cc}
Reply-To: {reply_to}
Subject: {subject}
Date: {iso_date}
Message-ID: {message['id']}
Labels: {label_str}
Attachments: {len(attachments)}

{body}
"""
                
                if attachments:
                    email_content += "\n\n--- Attachments ---\n"
                    for att in attachments:
                        status = "✓ Text Extracted" if att['text_extracted'] else "◯ Binary File"
                        email_content += f"- {att['filename']} ({att['size_mb']:.2f}MB) [{status}]\n"
                
                temp_filepath = filepath.with_suffix('.tmp')
                with open(temp_filepath, 'w', encoding='utf-8') as f:
                    f.write(email_content)
                
                temp_filepath.rename(filepath)
                
                save_synced_id(message['id'])
                synced_ids.add(message['id'])
                
                saved_count += 1
                att_info = f" ({len(attachments)} attachments)" if attachments else ""
                logger.info(f"✅ Saved: {subject[:60]}...{att_info}")
            
            except Exception as e:
                logger.error(f"Error processing {message['id']}: {e}")
                continue
        
        return saved_count
    
    except HttpError as error:
        logger.error(f'Gmail API error: {error}')
        return 0


def watch_gmail(service, check_interval=60, search_query="newer_than:7d"):
    """Continuously watch for new emails"""
    
    print("=" * 80)
    print("          📧 Gmail Watcher V2 - Enhanced Edition")
    print("=" * 80)
    print(f"✓ Monitoring Gmail inbox")
    print(f"✓ Check interval: {check_interval} seconds")
    print(f"✓ Search query: '{search_query or 'All emails'}'")
    print(f"✓ Output: {OUTPUT_DIR}")
    print(f"✓ Attachments: {ATTACHMENTS_DIR}")
    print(f"✓ Extracted text: {PROCESSED_DIR}")
    print(f"\n📎 Attachment Support: {'ENABLED' if ATTACHMENT_CONFIG['enabled'] else 'DISABLED'}")
    
    if ATTACHMENT_CONFIG['enabled']:
        print(f"  • Max file size: {ATTACHMENT_CONFIG['max_file_size_mb']}MB")
        print(f"  • PDF page limit: {ATTACHMENT_CONFIG['pdf_max_pages'] or 'None'}")
        print(f"  • Supported types ({len(ATTACHMENT_CONFIG['allowed_extensions'])}):")
        print(f"    Documents: PDF, DOCX, DOC, RTF, ODT, TXT")
        print(f"    Spreadsheets: CSV, XLSX, XLS, ODS")
        print(f"    Presentations: PPTX, PPT, ODP")
        print(f"    Images: JPG, PNG, GIF, BMP, TIFF, WEBP (with OCR)")
        print(f"    Web/Markup: HTML, XML, JSON, Markdown")
        print(f"    Code: Python, JavaScript, Java, C++, SQL, YAML, etc.")
        print(f"    Archives: ZIP, RAR, 7Z, TAR, GZ (metadata only)")
        print(f"  • 🚀 Instant text extraction: {ATTACHMENT_CONFIG['extract_text_immediately']}")
    
    print("\n⚠️  Keep this running alongside app.py for real-time updates")
    print("   Press Ctrl+C to stop\n")
    print("=" * 80)
    
    synced_ids = load_synced_ids()
    logger.info(f"📊 Already synced: {len(synced_ids)} emails")
    
    check_count = 0
    
    try:
        while True:
            check_count += 1
            timestamp = datetime.now().strftime("%H:%M:%S")
            
            logger.info(f"[{timestamp}] Check #{check_count} - Looking for new emails...")
            
            new_count = fetch_new_emails(
                service, 
                synced_ids, 
                search_query=search_query,
                max_check=200
            )
            
            if new_count > 0:
                logger.info(f"✅ Downloaded {new_count} new emails")
                logger.info("   Pathway will automatically index them!")
            else:
                logger.info("   No new emails")
            
            logger.info(f"   Next check in {check_interval} seconds...")
            time.sleep(check_interval)
    
    except KeyboardInterrupt:
        print("\n\n👋 Stopping Gmail watcher...")
        print(f"✅ Total emails synced: {len(synced_ids)}")
        print("=" * 80)


def initial_sync(service, max_emails=200, search_query="newer_than:30d"):
    """Do initial sync of recent emails"""
    
    logger.info(f"\n🔄 Running initial sync...")
    logger.info(f"   Fetching up to {max_emails} emails")
    logger.info(f"   Query: '{search_query}'")
    
    synced_ids = load_synced_ids()
    new_count = fetch_new_emails(
        service,
        synced_ids,
        search_query=search_query,
        max_check=min(max_emails, 200)  # Gmail API limit per request
    )
    
    logger.info(f"\n✅ Initial sync complete: {new_count} new emails downloaded")
    
    return synced_ids


def print_dependency_check():
    """Check and print status of optional dependencies"""
    
    print("\n" + "=" * 80)
    print("           📦 Dependency Status Check")
    print("=" * 80)
    
    deps = {
        'Required': {
            'google-auth': 'Gmail API authentication',
            'google-auth-oauthlib': 'Gmail OAuth',
            'google-api-python-client': 'Gmail API client',
        },
        'Core Features': {
            'pandas': 'CSV/Excel processing',
            'openpyxl': 'Excel (.xlsx) support',
            'PyPDF2': 'PDF text extraction',
            'python-docx': 'Word (.docx) support',
            'beautifulsoup4': 'HTML parsing',
        },
        'Enhanced Features': {
            'python-pptx': 'PowerPoint support',
            'pillow': 'Image handling',
            'pytesseract': 'Image OCR (also needs Tesseract)',
            'striprtf': 'RTF file support',
            'odfpy': 'OpenDocument support',
            'tabulate': 'Better table formatting',
        },
        'Optional': {
            'textract': 'Legacy .doc files',
            'rarfile': 'RAR archive support',
        }
    }
    
    for category, packages in deps.items():
        print(f"\n{category}:")
        for package, description in packages.items():
            try:
                __import__(package.replace('-', '_'))
                print(f"  ✅ {package:25} - {description}")
            except ImportError:
                print(f"  ❌ {package:25} - {description}")
    
    print("\n" + "=" * 80)
    print("\nTo install all dependencies:")
    print("pip install pandas openpyxl PyPDF2 python-docx python-pptx")
    print("pip install beautifulsoup4 pillow pytesseract striprtf odfpy tabulate")
    print("\nFor OCR: Also install Tesseract OCR engine")
    print("  https://github.com/tesseract-ocr/tesseract")
    print("=" * 80 + "\n")


def main():
    """Main function"""
    
    # Configuration
    CHECK_INTERVAL = 60
    SEARCH_QUERY = "newer_than:7d"
    INITIAL_SYNC_MAX = 200
    
    print_dependency_check()
    
    print("\n📧 Gmail Watcher V2 Configuration:")
    print(f"   Check interval: {CHECK_INTERVAL} seconds")
    print(f"   Search query: '{SEARCH_QUERY}'")
    print(f"   Initial sync max: {INITIAL_SYNC_MAX} emails")
    
    logger.info("\n🔐 Authenticating with Gmail...")
    creds = authenticate_gmail()
    
    if not creds:
        return
    
    service = build('gmail', 'v1', credentials=creds)
    logger.info("✅ Authentication successful!")
    
    synced_ids = initial_sync(service, INITIAL_SYNC_MAX, SEARCH_QUERY)
    
    print("\n" + "=" * 80)
    input("Press Enter to start real-time monitoring (or Ctrl+C to exit)...")
    
    watch_gmail(service, check_interval=CHECK_INTERVAL, search_query=SEARCH_QUERY)


if __name__ == "__main__":
    main()
